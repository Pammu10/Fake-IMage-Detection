import datetime
import json
import os
from typing import List

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


RESULTS_PATH = os.path.join("outputs", "results.json")
PPT_OUTPUT_PATH = os.path.join("outputs", "AI_vs_Real_Image_Detection.pptx")


BG_COLOR = RGBColor(20, 24, 32)
TITLE_COLOR = RGBColor(240, 240, 240)
TEXT_COLOR = RGBColor(210, 210, 210)
ACCENT_COLOR = RGBColor(83, 177, 219)
CARD_COLOR = RGBColor(34, 41, 54)
FONT_NAME = "Calibri"


def apply_dark_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_title(slide, text: str):
    shape = slide.shapes.title
    shape.text = text
    p = shape.text_frame.paragraphs[0]
    p.font.name = FONT_NAME
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = TITLE_COLOR


def add_bullets(slide, lines: List[str], left=0.7, top=1.6, width=12.0, height=4.8, font_size=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = FONT_NAME
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR


def add_interpretation(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.6))
    tf = box.text_frame
    tf.text = f"Interpretation: {text}"
    p = tf.paragraphs[0]
    p.font.name = FONT_NAME
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = ACCENT_COLOR


def add_footer(slide, text: str = "AI vs Real Image Detection"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 160, 170)
    p.alignment = PP_ALIGN.RIGHT


def add_metric_cards(slide, items):
    left = 0.9
    top = 1.8
    card_w = 3.7
    card_h = 1.35
    gap = 0.35

    for title, value in items:
        shape = slide.shapes.add_shape(
            1,
            Inches(left),
            Inches(top),
            Inches(card_w),
            Inches(card_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_COLOR
        shape.line.color.rgb = ACCENT_COLOR

        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_NAME
        p1.font.size = Pt(14)
        p1.font.color.rgb = TEXT_COLOR

        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.name = FONT_NAME
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TITLE_COLOR

        left += card_w + gap


def add_image(slide, image_path: str, left=0.9, top=1.4, max_width=11.4, max_height=4.9):
    if os.path.exists(image_path):
        with Image.open(image_path) as im:
            w, h = im.size
        aspect = w / max(h, 1)

        target_w = max_width
        target_h = target_w / max(aspect, 1e-8)
        if target_h > max_height:
            target_h = max_height
            target_w = target_h * aspect

        x = left + (max_width - target_w) / 2
        y = top + (max_height - target_h) / 2
        slide.shapes.add_picture(image_path, Inches(x), Inches(y), width=Inches(target_w), height=Inches(target_h))
    else:
        add_bullets(slide, [f"Image not found: {image_path}"], top=3.0, font_size=18)


def add_table(slide, rows: List[List[str]], left=0.8, top=1.8, width=11.6, height=4.5):
    n_rows = len(rows)
    n_cols = len(rows[0]) if rows else 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = str(rows[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_NAME
                p.font.size = Pt(12 if n_cols >= 6 else 14)
                p.font.color.rgb = TEXT_COLOR if r > 0 else TITLE_COLOR
                p.alignment = PP_ALIGN.CENTER


def add_highlight_box(slide, title: str, text: str, left=0.8, top=5.55, width=11.6, height=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_COLOR
    shape.line.color.rgb = ACCENT_COLOR

    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = FONT_NAME
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_COLOR

    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.name = FONT_NAME
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_COLOR


def main():
    if not os.path.exists(RESULTS_PATH):
        raise FileNotFoundError("results.json not found. Run train.py first.")

    with open(RESULTS_PATH, "r", encoding="utf-8") as f:
        results = json.load(f)

    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_dark_background(slide)
    add_title(slide, "AI vs Real Image Detection")
    subtitle = slide.placeholders[1]
    subtitle.text = f"Student Name: <Your Name>\nDate: {datetime.date.today()}"
    for p in subtitle.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
    add_footer(slide)

    # Slide 2: Problem statement
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Problem Statement")
    add_bullets(
        slide,
        [
            "Goal: classify images as real photos or AI-generated images.",
            "Why it matters: synthetic media is growing quickly and can spread misinformation.",
            "Challenge: AI images can look realistic, but often contain subtle global inconsistencies.",
        ],
    )
    add_footer(slide)

    # Slide 3: Dataset
    ds = results["dataset"]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Dataset: CIFAKE")
    add_bullets(
        slide,
        [
            f"Total samples: {ds['total_samples']}",
            f"Classes: {', '.join(ds['class_names'])} (balanced)",
            f"Split: Train={ds['split']['train']}, Val={ds['split']['val']}, Test={ds['split']['test']}",
            "Dataset source: Kaggle CIFAKE (real + AI-generated images).",
        ],
    )
    add_footer(slide)

    # Slide 4: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Architecture Chosen")
    add_bullets(
        slide,
        [
            "Input Image (224x224)",
            "-> Frozen ViT Backbone (google/vit-base-patch16-224)",
            "-> CLS Token (768)",
            "-> Linear(768->256) + ReLU + Dropout(0.3)",
            "-> Linear(256->2)",
            "Why ViT: captures global image structure with self-attention, useful for detecting subtle AI artifacts.",
        ],
        font_size=19,
    )
    add_footer(slide)

    # Slide 5: Training setup
    setup = results["training_setup"]
    sweep = results["hyperparameter_sweep"]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Training Setup")
    lines = [
        f"Loss: {setup['loss']}",
        f"Optimizer: {setup['optimizer']}",
        f"Scheduler: {setup['scheduler']}",
        "Hyperparameters tried:",
    ]
    for row in sweep:
        lines.append(
            f"- {row['Config']}: lr={row['Learning Rate']}, batch={row['Batch Size']}, epochs={row['Epochs']}"
        )
    add_bullets(slide, lines, font_size=19)
    add_footer(slide)

    # Slide 6: Hyperparameter comparison table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Hyperparameter Comparison")
    header = ["Config", "LR", "Batch", "Epochs", "Val Acc", "Val Loss", "Val AUC"]
    rows = [header]
    for row in sweep:
        rows.append(
            [
                row["Config"],
                row["Learning Rate"],
                row["Batch Size"],
                row["Epochs"],
                row["Best Val Acc"],
                row["Best Val Loss"],
                row["Best Val AUC"],
            ]
        )
    add_table(slide, rows)
    add_interpretation(slide, results["interpretations"].get("hyperparameter_table", ""))
    add_footer(slide)

    # Slide 7: Training curves
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Training and Validation Curves")
    add_image(slide, results["graphs"]["training_curves"])
    add_interpretation(slide, results["interpretations"].get("training_curves", ""))
    add_footer(slide)

    # Slide 8: Confusion matrix
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Confusion Matrix")
    add_image(slide, results["graphs"]["confusion_matrix"])
    add_interpretation(slide, results["interpretations"].get("confusion_matrix", ""))
    add_footer(slide)

    # Slide 9: ROC curve
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "ROC Curve")
    add_image(slide, results["graphs"]["roc_curve"])
    add_interpretation(slide, results["interpretations"].get("roc_curve", ""))
    add_footer(slide)

    # Slide 10: Precision-Recall curve
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Precision-Recall Curve")
    add_image(slide, results["graphs"]["precision_recall_curve"])
    add_interpretation(slide, results["interpretations"].get("precision_recall_curve", ""))
    add_footer(slide)

    # Slide 11: Confidence distribution
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Confidence Score Distribution")
    add_image(slide, results["graphs"]["confidence_distribution"])
    add_interpretation(slide, results["interpretations"].get("confidence_distribution", ""))
    add_footer(slide)

    # Slide 12: K-Fold results
    kfold = results["kfold_results"]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "3-Fold Cross Validation Results")
    k_rows = [["Fold", "Val Accuracy", "Val AUC"]]
    for fold in kfold["folds"]:
        k_rows.append([fold["fold"], round(fold["val_accuracy"], 4), round(fold["val_auc"], 4)])
    k_rows.append(
        [
            "Mean±Std",
            f"{kfold['mean_accuracy']:.4f} ± {kfold['std_accuracy']:.4f}",
            f"{kfold['mean_auc']:.4f} ± {kfold['std_auc']:.4f}",
        ]
    )
    add_table(slide, k_rows)
    kfold_easy = results.get("simple_explanation", {}).get(
        "kfold",
        "K-Fold means train multiple times with different validation splits to check consistency.",
    )
    add_interpretation(slide, f"{kfold_easy} {results['interpretations'].get('kfold_table', '')}")
    add_footer(slide)

    # Slide 13: Key findings and conclusion
    best = results["best_config"]
    final_metrics = results["final_test_metrics"]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "Key Findings and Conclusion")
    add_bullets(
        slide,
        [
            f"Best config: {best['name']} (lr={best['lr']}, batch={best['batch_size']}, epochs={best['epochs']})",
            f"Final Test Accuracy: {final_metrics['accuracy']:.4f}",
            f"Final Test AUC: {final_metrics['auc']:.4f}",
            "What worked: frozen pretrained ViT + lightweight head gave stable performance.",
            "Future improvements: unfreeze last ViT block, stronger augmentations, test more transformers.",
        ],
        top=3.35,
        height=2.1,
        font_size=18,
    )
    add_metric_cards(
        slide,
        [
            ("Best Config", str(best["name"])),
            ("Test Accuracy", f"{final_metrics['accuracy']:.3f}"),
            ("Test AUC", f"{final_metrics['auc']:.3f}"),
        ],
    )
    add_highlight_box(
        slide,
        "One-line takeaway",
        "Config A gave the strongest validation performance and stable K-Fold behavior, so it is selected as the final model.",
    )
    add_footer(slide)

    # Slide 14: References
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_dark_background(slide)
    add_title(slide, "References")
    add_bullets(
        slide,
        [
            "CNNDetect — Wang et al., CVPR 2020",
            "UniversalFakeDetect — Ojha et al., CVPR 2023",
            "DIRE — Wang et al., ICCV 2023",
            "FreqNet — Tan et al., AAAI 2024",
        ],
        font_size=22,
    )
    add_footer(slide)

    os.makedirs("outputs", exist_ok=True)
    prs.save(PPT_OUTPUT_PATH)
    print(f"Presentation generated at: {PPT_OUTPUT_PATH}")


if __name__ == "__main__":
    main()
